import logging
from io import BytesIO

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Emu

logger = logging.getLogger(__name__)

SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
SLIDE_W_EMU = 9144000
SLIDE_H_EMU = 5143500

HIDE_NAV_SCRIPT = """() => {
    ['#nav-prev','#nav-next','#nav-dots','#nav-progress','#nav-counter']
        .forEach(sel => { const e = document.querySelector(sel); if(e) e.style.display='none'; });
    const s = document.createElement('style');
    s.textContent = '*, *::before, *::after { animation-play-state: paused !important; transition: none !important; }';
    document.head.appendChild(s);
}"""

SHOW_SLIDE_SCRIPT = """(idx) => {
    const deck = document.getElementById('deck');
    if (deck) {
        deck.style.transition = 'none';
        deck.style.transform = 'translateX(' + (-idx * 1920) + 'px)';
    }
    document.querySelectorAll('section.slide, div.slide').forEach((s, j) => {
        if (j === idx) {
            s.classList.add('active');
            s.classList.remove('hidden','inactive','prev','next');
            s.style.visibility = 'visible';
            s.style.opacity = '1';
            s.style.zIndex = '10';
            s.style.display = '';
        } else {
            s.classList.remove('active');
            s.style.visibility = 'hidden';
            s.style.opacity = '0';
            s.style.zIndex = '0';
        }
    });
}"""


async def generate_pptx(html_content: str) -> bytes:
    async with async_playwright() as pw:
        browser = await pw.chromium.launch(
            headless=True,
            args=["--no-sandbox", "--disable-setuid-sandbox"],
        )
        context = await browser.new_context(
            viewport={"width": SLIDE_W_PX, "height": SLIDE_H_PX}
        )
        page = await context.new_page()
        await page.set_content(html_content, timeout=30000)
        await page.wait_for_load_state("networkidle", timeout=30000)
        await page.evaluate(HIDE_NAV_SCRIPT)

        slide_count = await page.locator("section.slide, div.slide").count()
        logger.info("PPTX 导出：检测到 %d 张幻灯片，开始截图", slide_count)
        screenshots: list[bytes] = []

        for i in range(slide_count):
            await page.evaluate(SHOW_SLIDE_SCRIPT, i)
            await page.wait_for_timeout(200)

            locator = page.locator("section.slide, div.slide").nth(i)
            bbox = await locator.bounding_box()
            if bbox and bbox["width"] > 100:
                shot = await page.screenshot(clip={
                    "x": bbox["x"],
                    "y": bbox["y"],
                    "width": min(bbox["width"], SLIDE_W_PX),
                    "height": min(bbox["height"], SLIDE_H_PX),
                })
            else:
                shot = await page.screenshot(full_page=False)

            screenshots.append(shot)
            logger.info("PPTX 第 %d/%d 张截图完成: %d KB", i + 1, slide_count, len(shot) // 1024)

        await browser.close()

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    blank_layout = prs.slide_layouts[6]

    for shot_bytes in screenshots:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            BytesIO(shot_bytes), left=0, top=0,
            width=Emu(SLIDE_W_EMU), height=Emu(SLIDE_H_EMU),
        )

    output = BytesIO()
    prs.save(output)
    pptx_bytes = output.getvalue()
    logger.info("PPTX 生成完成: %d 张幻灯片, 文件大小 %d KB", len(screenshots), len(pptx_bytes) // 1024)
    return pptx_bytes
