"""Tests for PPTX export — unit (service layer) and integration (API layer)."""
import io
import struct
import zlib
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from pptx import Presentation

from src.services.pptx_exporter import (
    SLIDE_W_EMU,
    SLIDE_H_EMU,
    HIDE_NAV_SCRIPT,
    SHOW_SLIDE_SCRIPT,
)

FAKE_PPTX = b"PK\x03\x04 fake pptx content"


def _make_valid_png() -> bytes:
    """Return a minimal valid 1x1 red PNG image (stdlib only)."""

    def _chunk(ctype: bytes, data: bytes) -> bytes:
        c = ctype + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)  # 1x1, 8-bit RGB
    idat = zlib.compress(b"\x00\xff\x00\x00")  # filter byte + R, G, B
    return b"\x89PNG\r\n\x1a\n" + _chunk(b"IHDR", ihdr) + _chunk(b"IDAT", idat) + _chunk(b"IEND", b"")


_VALID_PNG = _make_valid_png()


# ─── helpers ────────────────────────────────────────────────────────


def _make_fake_screenshot(size: int = 4096) -> bytes:
    return _VALID_PNG


def _build_mock_playwright_chain(slide_count: int = 2, *, bbox_override=None):
    """Build a complete mock Playwright chain for unit testing generate_pptx.

    Returns a tuple of (async_playwright_mock, mock_browser, mock_page) so
    callers can add extra assertions.
    """
    mock_page = MagicMock()
    mock_page.set_content = AsyncMock()
    mock_page.wait_for_load_state = AsyncMock()
    mock_page.wait_for_timeout = AsyncMock()
    mock_page.evaluate = AsyncMock()
    mock_page.screenshot = AsyncMock(return_value=_make_fake_screenshot())

    mock_locator = MagicMock()
    mock_locator.count = AsyncMock(return_value=slide_count)

    mock_nth_locator = MagicMock()
    mock_nth_locator.bounding_box = AsyncMock(
        return_value=bbox_override or {"x": 0, "y": 0, "width": 1920, "height": 1080}
    )

    mock_page.locator = MagicMock(return_value=mock_locator)
    mock_locator.nth = MagicMock(return_value=mock_nth_locator)

    mock_browser = MagicMock()
    mock_browser.close = AsyncMock()
    mock_browser.new_context = AsyncMock()
    mock_browser.new_context.return_value.new_page = AsyncMock(return_value=mock_page)

    mock_pw = MagicMock()
    mock_pw.chromium.launch = AsyncMock(return_value=mock_browser)
    mock_pw.__aenter__ = AsyncMock(return_value=mock_pw)
    mock_pw.__aexit__ = AsyncMock(return_value=None)

    mock_async_playwright = MagicMock()
    mock_async_playwright.__aenter__ = AsyncMock(return_value=mock_pw)
    mock_async_playwright.__aexit__ = AsyncMock(return_value=None)

    return mock_async_playwright, mock_browser, mock_page


# ─── unit tests: generate_pptx ──────────────────────────────────────


class TestGeneratePptx:
    """Unit tests for the generate_pptx service function."""

    @pytest.mark.asyncio
    async def test_generates_valid_pptx_from_single_slide(self):
        mock_ap, _, _ = _build_mock_playwright_chain(slide_count=1)
        html = "<section class='slide'>Hello</section>"

        with patch("src.services.pptx_exporter.async_playwright", return_value=mock_ap):
            from src.services.pptx_exporter import generate_pptx
            result = await generate_pptx(html)

        # Result is a valid OOXML ZIP
        assert result[:2] == b"PK"
        # Verify it's a valid pptx by re-reading it
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 1
        assert prs.slide_width == SLIDE_W_EMU
        assert prs.slide_height == SLIDE_H_EMU

    @pytest.mark.asyncio
    async def test_generates_pptx_with_multiple_slides(self):
        mock_ap, _, _ = _build_mock_playwright_chain(slide_count=3)
        html = "<div class='slide'>1</div><div class='slide'>2</div><div class='slide'>3</div>"

        with patch("src.services.pptx_exporter.async_playwright", return_value=mock_ap):
            from src.services.pptx_exporter import generate_pptx
            result = await generate_pptx(html)

        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 3

    @pytest.mark.asyncio
    async def test_empty_html_returns_pptx_with_zero_slides(self):
        mock_ap, _, _ = _build_mock_playwright_chain(slide_count=0)
        html = "<div>no slides here</div>"

        with patch("src.services.pptx_exporter.async_playwright", return_value=mock_ap):
            from src.services.pptx_exporter import generate_pptx
            result = await generate_pptx(html)

        assert result[:2] == b"PK"
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 0

    @pytest.mark.asyncio
    async def test_injects_hide_nav_script(self):
        mock_ap, _, mock_page = _build_mock_playwright_chain(slide_count=1)
        html = "<section class='slide'>x</section>"

        with patch("src.services.pptx_exporter.async_playwright", return_value=mock_ap):
            from src.services.pptx_exporter import generate_pptx
            await generate_pptx(html)

        mock_page.evaluate.assert_any_call(HIDE_NAV_SCRIPT)

    @pytest.mark.asyncio
    async def test_injects_show_slide_script_per_slide(self):
        slide_count = 2
        mock_ap, _, mock_page = _build_mock_playwright_chain(slide_count=slide_count)
        html = "<section class='slide'>A</section><section class='slide'>B</section>"

        with patch("src.services.pptx_exporter.async_playwright", return_value=mock_ap):
            from src.services.pptx_exporter import generate_pptx
            await generate_pptx(html)

        # SHOW_SLIDE_SCRIPT called with idx 0 and idx 1
        show_calls = [
            c for c in mock_page.evaluate.call_args_list
            if c.args and c.args[0] == SHOW_SLIDE_SCRIPT
        ]
        assert len(show_calls) == slide_count
        assert show_calls[0].args[1] == 0
        assert show_calls[1].args[1] == 1

    @pytest.mark.asyncio
    async def test_falls_back_to_fullpage_when_bbox_too_small(self):
        mock_ap, _, mock_page = _build_mock_playwright_chain(
            slide_count=1,
            bbox_override={"x": 0, "y": 0, "width": 50, "height": 50},
        )
        html = "<div class='slide'>x</div>"

        with patch("src.services.pptx_exporter.async_playwright", return_value=mock_ap):
            from src.services.pptx_exporter import generate_pptx
            await generate_pptx(html)

        # Should have used full_page=False screenshot
        mock_page.screenshot.assert_called_with(full_page=False)

    @pytest.mark.asyncio
    async def test_screenshots_each_slide_with_correct_clip(self):
        mock_ap, _, mock_page = _build_mock_playwright_chain(slide_count=2)
        html = "<section class='slide'>A</section><section class='slide'>B</section>"

        with patch("src.services.pptx_exporter.async_playwright", return_value=mock_ap):
            from src.services.pptx_exporter import generate_pptx
            await generate_pptx(html)

        assert mock_page.screenshot.call_count == 2
        # Each call should use clip with the bounding box coordinates
        for call_args in mock_page.screenshot.call_args_list:
            assert "clip" in call_args.kwargs
            assert call_args.kwargs["clip"]["width"] == 1920
            assert call_args.kwargs["clip"]["height"] == 1080


# ─── integration tests: API endpoint ────────────────────────────────


class TestExportPptx:
    """API-level integration tests for the PPTX export endpoint."""

    def test_export_pptx_returns_pptx(self, api_client, temp_workspace):
        create_resp = api_client.post("/api/v1/projects", json={"name": "Deck"})
        pid = create_resp.json()["id"]
        proj_dir = temp_workspace / "projects" / pid
        (proj_dir / "stub.html").write_text("")

        html = b"<html><style>.slide{}</style><div class='slide'>x</div></html>"
        upload_resp = api_client.post(
            "/api/v1/files/upload",
            files={"file": ("deck.html", io.BytesIO(html), "text/html")},
            data={"project_id": pid},
        )
        assert upload_resp.status_code == 200

        with patch(
            "src.api.export.generate_pptx", new_callable=AsyncMock
        ) as mock_gen:
            mock_gen.return_value = FAKE_PPTX
            resp = api_client.get(f"/api/v1/projects/{pid}/export/pptx")

        assert resp.status_code == 200
        assert resp.headers["content-type"] == (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        assert "attachment" in resp.headers["content-disposition"]
        assert resp.content == FAKE_PPTX
        mock_gen.assert_awaited_once()

    def test_export_pptx_empty_project_returns_404(self, api_client):
        create_resp = api_client.post("/api/v1/projects", json={"name": "Empty"})
        pid = create_resp.json()["id"]

        resp = api_client.get(f"/api/v1/projects/{pid}/export/pptx")
        assert resp.status_code == 404

    def test_export_pptx_nonexistent_project_returns_404(self, api_client):
        resp = api_client.get("/api/v1/projects/proj-000000000000/export/pptx")
        assert resp.status_code == 404

    def test_export_pptx_generation_error_returns_500(self, api_client, temp_workspace):
        create_resp = api_client.post("/api/v1/projects", json={"name": "Deck"})
        pid = create_resp.json()["id"]
        proj_dir = temp_workspace / "projects" / pid
        (proj_dir / "stub.html").write_text("")

        html = b"<html><div class='slide'>x</div></html>"
        upload_resp = api_client.post(
            "/api/v1/files/upload",
            files={"file": ("deck.html", io.BytesIO(html), "text/html")},
            data={"project_id": pid},
        )
        assert upload_resp.status_code == 200

        with patch(
            "src.api.export.generate_pptx", new_callable=AsyncMock
        ) as mock_gen:
            mock_gen.side_effect = RuntimeError("Playwright crashed")
            resp = api_client.get(f"/api/v1/projects/{pid}/export/pptx")

        assert resp.status_code == 500

    def test_export_pptx_filename_in_content_disposition(self, api_client, temp_workspace):
        create_resp = api_client.post("/api/v1/projects", json={"name": "My Deck"})
        pid = create_resp.json()["id"]
        proj_dir = temp_workspace / "projects" / pid
        (proj_dir / "stub.html").write_text("")

        html = b"<html><div class='slide'>x</div></html>"
        api_client.post(
            "/api/v1/files/upload",
            files={"file": ("deck.html", io.BytesIO(html), "text/html")},
            data={"project_id": pid},
        )

        with patch(
            "src.api.export.generate_pptx", new_callable=AsyncMock
        ) as mock_gen:
            mock_gen.return_value = FAKE_PPTX
            resp = api_client.get(f"/api/v1/projects/{pid}/export/pptx")

        assert resp.status_code == 200
        cd = resp.headers["content-disposition"]
        assert "My-Deck.pptx" in cd

    def test_export_pptx_passes_html_to_generator(self, api_client, temp_workspace):
        create_resp = api_client.post("/api/v1/projects", json={"name": "Deck"})
        pid = create_resp.json()["id"]
        proj_dir = temp_workspace / "projects" / pid
        (proj_dir / "stub.html").write_text("")

        html = b"<html><body><div class='slide'>Hello World</div></body></html>"
        api_client.post(
            "/api/v1/files/upload",
            files={"file": ("deck.html", io.BytesIO(html), "text/html")},
            data={"project_id": pid},
        )

        with patch(
            "src.api.export.generate_pptx", new_callable=AsyncMock
        ) as mock_gen:
            mock_gen.return_value = FAKE_PPTX
            api_client.get(f"/api/v1/projects/{pid}/export/pptx")

        # Verify the HTML content was passed to generate_pptx
        called_html = mock_gen.call_args[0][0]
        assert "Hello World" in called_html
